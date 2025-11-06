import os, re, json, docx, openpyxl, camelot, asyncio, sys, multiprocessing as multi
from langchain_community.vectorstores import FAISS;from langchain_community.embeddings import HuggingFaceEmbeddings as Embed;from langchain_text_splitters import RecursiveCharacterTextSplitter as Rec
from mcp.server import Server ; from mcp.server.stdio import stdio_server ;from mcp.types import Tool, TextContent
from datetime import datetime
from pptx import Presentation
from openpyxl.utils import column_index_from_string
from PyPDF2 import PdfReader
from concurrent.futures import ThreadPoolExecutor

app = Server("excel-mcp"); batch_size = 64
embeddings = Embed(model_name="BAAI/bge-small-en-v1.5", model_kwargs={"device":"cuda"}, encode_kwargs={"normalize_embeddings":True,"batch_size":batch_size})
FAISS.from_texts(["warmup"],embeddings); print("FAISS loaded", file=sys.stderr, flush=True)

def _is_numbery(s): return s.replace(",","").replace("$","").replace(".","").replace("-","").isdigit()
def _excel_cell_to_str(c, include_formulas):
    if c.value is None: return ""
    if getattr(c, "data_type", None) == "f":
        if include_formulas and hasattr(c, "formula"): return f"{c.value} [{c.formula}]"
        if include_formulas: return f"{c.value} [{c.value}]"
        return str(c.value)
    if isinstance(c.value, (int, float)): return str(round(c.value, 2))
    if isinstance(c.value, datetime): return c.value.date().isoformat()
    return str(c.value)

def parse_file(path, include_formulas=False, chunk_size=4000, chunk_overlap=100):
    full_docs=[]; chunks=[]; metas=[]
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext in [".xls", ".xlsx"]:
            wb = openpyxl.load_workbook(path, data_only=not include_formulas)
            for ws in wb.worksheets:
                rows=[]; row_nums=[]
                for r in ws.iter_rows(values_only=False):
                    vals=[_excel_cell_to_str(c, include_formulas) for c in r]
                    while vals and (not vals[-1].strip() or vals[-1].strip()=="0"): vals.pop()
                    if any(vals): rows.append("|".join(vals)); row_nums.append(r[0].row)
                if rows:
                    full = f"\n[{path} | Sheet: {ws.title}]\n" + "\n".join(rows); full_docs.append(full)
                    cur,cur_rows,cl=[],[],0
                    for rn,rt in zip(row_nums,rows):
                        l=len(rt)+1
                        if cl+l>chunk_size and cur:
                            rr=f"{cur_rows[0]}-{cur_rows[-1]}"; t=f"[{path}|Sheet:{ws.title}|Rows:{rr}]\n"+'\n'.join(cur)
                            chunks.append(t); metas.append({"file":path,"pg":ws.title,"rows":rr,"text":'\n'.join(cur),"chunk_type":"sheet"}); cur,cur_rows,cl=[],[],0
                        cur.append(rt); cur_rows.append(rn); cl+=l
                    if cur:
                        rr=f"{cur_rows[0]}-{cur_rows[-1]}"; t=f"[{path}|Sheet:{ws.title}|Rows:{rr}]\n"+'\n'.join(cur)
                        chunks.append(t); metas.append({"file":path,"pg":ws.title,"rows":rr,"text":'\n'.join(cur),"chunk_type":"sheet"})
        elif ext == ".pdf":
            reader = PdfReader(path)
            for i in range(1, len(reader.pages)+1):
                lines=[line.strip() for line in (reader.pages[i-1].extract_text() or "").splitlines() if line.strip()]
                num_cells=sum(1 for l in lines for c in re.split(r'\s{2,}',l) if _is_numbery(c))
                total_cells=sum(1 for l in lines for c in re.split(r'\s{2,}',l) if c)
                use_camelot = total_cells>0 and (num_cells/total_cells)>=0.2
                if use_camelot:
                    try:
                        tabrows=["|".join(str(c).strip() or "-" for c in r) for r in camelot.read_pdf(path,pages=str(i),flavor="stream",strip_text="\n",edge_tol=500)[0].df.values]
                        rows = tabrows
                    except: rows = ["|".join(c for c in re.split(r'\s{2,}',l)) for l in lines]
                else: rows = ["|".join(c for c in re.split(r'\s{2,}',l)) for l in lines]
                if rows:
                    full_docs.append(f"\n[{path} | Page {i}]\n" + "\n".join(rows))
                    cur,ch_rows,cl,residual_pg=[],[],0,[]
                    for rn,r in enumerate(rows,1):
                        l=len(r)+1
                        if cl+l>chunk_size and cur:
                            rr=f"{ch_rows[0]}-{ch_rows[-1]}"; pg_label=','.join(residual_pg) or str(i)
                            chunks.append((f"[{path}|Pg:{pg_label}|Rows:{rr}]\n"+'\n'.join(cur))); metas.append({"file":path,"pg":pg_label,"rows":rr,"chunk_type":"pdf_table" if use_camelot else "pdf_text","text":'\n'.join(cur)})
                            cur,ch_rows,cl,residual_pg=[],[],0,[]
                        cur.append(r); ch_rows.append(rn); cl+=l; residual_pg.append(str(i))
                    if cur:
                        rr=f"{ch_rows[0]}-{ch_rows[-1]}"; chunks.append((f"[{path}|Pg:{i}|Rows:{rr}]\n"+'\n'.join(cur))); metas.append({"file":path,"pg":i,"rows":rr,"chunk_type":"pdf_table" if use_camelot else "pdf_text","text":'\n'.join(cur)})
        elif ext in [".ppt",".pptx"]:
            prs = Presentation(path)
            content = "\n\n".join("\n".join(sh.text for sh in sl.shapes if hasattr(sh, "text")) for sl in prs.slides)
            if content.strip(): full_docs.append(f"\n[{path}]\n{content}")
            raw = Rec(chunk_size=chunk_size,chunk_overlap=chunk_overlap).split_text(content) if content.strip() else []
            for i,c in enumerate(raw): chunks.append(f"[{path}|Chunk:{i}]{c}"); metas.append({"file":path,"rows":i,"text":c,"chunk_type":"text_split"})
            while chunks and (not chunks[-1].strip() or len(chunks[-1].split())<5): chunks.pop(); metas.pop()
        elif ext in [".doc", ".docx"]:
            txt = "\n".join(p.text for p in docx.Document(path).paragraphs)
            if txt.strip(): full_docs.append(f"\n[{path}]\n{txt}")
            raw = Rec(chunk_size=chunk_size,chunk_overlap=chunk_overlap).split_text(txt) if txt.strip() else []
            for i,c in enumerate(raw): chunks.append(f"[{path}|Chunk:{i}]{c}"); metas.append({"file":path,"rows":i,"text":c,"chunk_type":"text_split"})
            while chunks and (not chunks[-1].strip() or len(chunks[-1].split())<5): chunks.pop(); metas.pop()
    except Exception as e:
        return {"path": path, "full_docs": full_docs, "chunks": [], "metas": []}
    return {"path": path, "full_docs": full_docs, "chunks": [c if isinstance(c,str) else c for c in chunks], "metas": metas}

@app.list_tools()
async def list_tools():
    return [Tool(name="query_docs", description="Extract detailed text and table content from Excel, Word, PDF, or PowerPoint files for context. IMPORTANT: Always pass ALL file paths in a SINGLE call - do not make multiple separate calls for different files.",
             inputSchema={"type":"object","properties":{"paths":{"type":"array","items":{"type":"string"}},"include_formulas":{"type":"boolean","description":"For Excel files, include formulas alongside values","default":False}},"required":["paths"]}),
            Tool(name="apply_excel_updates", description="""Apply structured Excel updates given JSON-formatted proposed changes.
            Expected JSON format (can be a single object or array of objects): {"Sheet": "SheetName","StartCell": "A1","Data": [["row1col1", ...]]}""",
             inputSchema={"type":"object","properties":{"path":{"type":"string"},"proposal":{"type":"string"}},"required":["path","proposal"]}),
            Tool(name="vectorize_embed", description="Vectorize folder contents or query vectors (loads .faiss if present).",
             inputSchema={"type":"object","properties":{"folder":{"type":"string"},"query":{"type":"string"},"k":{"type":"integer","default":10},"include_formulas":{"type":"boolean","default":False}},"required":["folder"]})
            ]

@app.call_tool()
async def call_tool(name: str, args: dict):
    if name == "query_docs":
        results=[]; include_formulas=args.get("include_formulas", False)
        for p in args["paths"]:
            r=parse_file(p, include_formulas=include_formulas)
            if r["full_docs"]: results.extend(r["full_docs"])
        all_chunks=[]; all_meta=[]
        for p in args["paths"]:
            r=parse_file(p, include_formulas=include_formulas)
            all_chunks.extend(r["chunks"]); all_meta.extend(r["metas"])
        payload={"full_docs":"\n".join(results) if results else "No content extracted","chunks":all_chunks,"metas":all_meta}
        return [TextContent(type="text", text=json.dumps(payload, indent=2))]
    if name == "apply_excel_updates":
        proposal=json.loads(args["proposal"]); wb=openpyxl.load_workbook(args["path"])
        updates=proposal if isinstance(proposal, list) else [proposal]
        for u in updates:
            sn=u.get("Sheet")
            if sn not in wb.sheetnames: continue
            ws=wb[sn]; sc=u.get("StartCell","A1")
            cm=re.match(r"([A-Z]+)", sc); rm=re.search(r"(\d+)", sc)
            if cm and rm:
                sc_idx=column_index_from_string(cm.group()); sr=int(rm.group())
                for i,row in enumerate(u.get("Data",[])):
                    for j,val in enumerate(row): ws.cell(row=sr+i, column=sc_idx+j, value=val)
        wb.save(args["path"]); return [TextContent(type="text", text=f"✅ Applied updates to {args['path']}")]
    if name == "vectorize_embed":
        folder=args["folder"]; query=args.get("query"); k=args.get("k",10); include_formulas=args.get("include_formulas", False)
        v_path=os.path.join(folder, ".faiss"); vector = FAISS.load_local(v_path,embeddings,allow_dangerous_deserialization=True) if os.path.exists(v_path) else None
        docs=[]
        for root,dirs,files in os.walk(folder):
            for f in files:
                if os.path.splitext(f)[1].lower() in [".xls",".xlsx",".pdf",".ppt",".pptx",".doc",".docx"]: docs.append(os.path.join(root,f))
        if query and not vector:
            all_chunks=[]; all_meta=[]
            with ThreadPoolExecutor(max_workers=min(20,multi.cpu_count()*2)) as ex:
                for r in ex.map(lambda p: parse_file(p, include_formulas=include_formulas), docs):
                    if not r["chunks"]: continue
                    all_chunks.extend(r["chunks"]); all_meta.extend(r["metas"])
            if all_chunks:
                vector=FAISS.from_texts(all_chunks,embeddings,metadatas=all_meta)
                vector.save_local(v_path)
        if query:
            if not vector: return [TextContent(type="text", text="No vector store found. Process files first.")]
            results=[{"content":r.page_content,"metadata":r.metadata or {}} for r in vector.similarity_search(query,k=k)]
            return [TextContent(type="text", text=json.dumps(results, indent=2))]
        existing = {d.metadata.get("file") for d in vector.docstore._dict.values() if d.metadata} if vector else set()
        all_chunks=[]; all_meta=[]
        max_workers=min(20,multi.cpu_count()*2)
        with ThreadPoolExecutor(max_workers=max_workers) as ex:
            for r in ex.map(lambda p: parse_file(p, include_formulas=include_formulas), [f for f in docs if f not in existing]):
                if not r["chunks"]: continue
                all_chunks.extend(r["chunks"]); all_meta.extend(r["metas"])
        for i in range(0, len(all_chunks), batch_size*4):
            batch_texts=all_chunks[i:i+batch_size*4]; batch_meta=all_meta[i:i+batch_size*4]
            if vector is None: vector=FAISS.from_texts(batch_texts, embeddings, metadatas=batch_meta)
            else: vector.add_texts(batch_texts, metadatas=batch_meta)
        if vector: vector.save_local(v_path)
        return [TextContent(type="text", text=f"✅ Vectorized {len(all_chunks)} chunks from {len(docs)} files")]

async def main():
    async with stdio_server() as (read_stream, write_stream): await app.run(read_stream, write_stream, app.create_initialization_options())
if __name__ == "__main__": asyncio.run(main())
